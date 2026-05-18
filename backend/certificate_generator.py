"""
Generates a personalized PDF certificate from a PPTX template.

Workflow:
  1. Open the PPTX template with python-pptx.
  2. Replace all {{ColumnName}} placeholders with values from the CSV row.
     - Placeholder format is EXACT and case-sensitive: {{Name}}, {{Email}}, etc.
     - The column name in the CSV must match the placeholder exactly.
  3. Save as a temporary PPTX.
  4. Convert to PDF via MS PowerPoint COM automation (requires PowerPoint installed).
  5. Delete the temporary PPTX.
  6. Return the path to the generated PDF.

Caller is responsible for deleting the returned PDF after sending.
"""

import hashlib
import os
import re
import subprocess
import sys
import time

def _replace_in_paragraph(paragraph, replacements: dict):
    """
    Replaces {{Key}} placeholders in a paragraph.

    PowerPoint sometimes splits visually-uniform text across multiple runs
    (e.g. {{Na | me}}). We join all runs, do the replacement, then write
    the result back into the first run and blank the rest — preserving the
    first run's character formatting (font, size, colour).
    """
    full_text = "".join(run.text for run in paragraph.runs)

    if not any(f"{{{{{k}}}}}" in full_text for k in replacements):
        return

    for key, value in replacements.items():
        full_text = full_text.replace(f"{{{{{key}}}}}", value)

    if paragraph.runs:
        paragraph.runs[0].text = full_text
        for run in paragraph.runs[1:]:
            run.text = ""


def _apply_replacements(template_path: str, output_path: str, replacements: dict):
    from pptx import Presentation

    prs = Presentation(template_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                _replace_in_paragraph(paragraph, replacements)
    prs.save(output_path)


_PP_SAVE_AS_PDF = 32  # ppSaveAsPDF constant


def _convert_to_pdf_win32(pptx_path: str, pdf_path: str):
    import pythoncom
    import win32com.client

    abs_pptx = os.path.abspath(pptx_path)
    abs_pdf  = os.path.abspath(pdf_path)

    pythoncom.CoInitialize()
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = True

    presentation = None
    try:
        presentation = powerpoint.Presentations.Open(abs_pptx, WithWindow=False)
        presentation.SaveAs(abs_pdf, _PP_SAVE_AS_PDF)
    finally:
        if presentation:
            presentation.Close()
        powerpoint.Quit()
        pythoncom.CoUninitialize()
        # Give Windows time to release the file handle before deleting the temp PPTX
        time.sleep(0.8)


def _convert_to_pdf_libreoffice(pptx_path: str, pdf_path: str):
    abs_pptx = os.path.abspath(pptx_path)
    out_dir  = os.path.dirname(abs_pptx)

    result = subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", out_dir, abs_pptx],
        capture_output=True,
        text=True,
        timeout=120,
    )
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice conversion failed: {result.stderr.strip()}")

    # LibreOffice outputs {stem}.pdf alongside the input file — rename to expected path
    lo_output = os.path.join(out_dir, os.path.splitext(os.path.basename(abs_pptx))[0] + ".pdf")
    abs_pdf   = os.path.abspath(pdf_path)
    if lo_output != abs_pdf:
        os.replace(lo_output, abs_pdf)


def _convert_to_pdf(pptx_path: str, pdf_path: str):
    if sys.platform == "win32":
        _convert_to_pdf_win32(pptx_path, pdf_path)
    else:
        _convert_to_pdf_libreoffice(pptx_path, pdf_path)


def _delete_retry(path: str, retries: int = 6, delay: float = 0.5):
    """
    Delete a file, retrying on PermissionError (Windows file-lock race).
    Never raises — a leftover temp file must not abort the send.
    """
    for attempt in range(retries):
        try:
            if os.path.exists(path):
                os.remove(path)
            return
        except PermissionError:
            if attempt < retries - 1:
                time.sleep(delay)
    # Still locked after all retries — leave it, do not raise


def generate(template_path: str, output_dir: str, row: dict) -> str:
    """
    Full pipeline: replace placeholders → convert to PDF → return PDF path.

    Args:
        template_path: Absolute or relative path to the .pptx template file.
        output_dir:    Directory where the temporary PPTX and final PDF are saved.
        row:           Dict of CSV column values (original case preserved).
                       Keys must match {{Placeholder}} names in the PPTX exactly.

    Returns:
        Absolute path to the generated PDF. Caller must delete it after use.
    """
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"PPTX template not found: {template_path}")

    os.makedirs(output_dir, exist_ok=True)

    name_val = next((v for k, v in row.items() if k.lower() == "name"), "certificate")
    email_val = next((v for k, v in row.items() if k.lower() == "email"), "")
    safe_name = re.sub(r"[^\w\s-]", "", name_val).strip().replace(" ", "_")
    # 8-char email hash makes filenames unique even when two participants share a name
    uid = hashlib.md5(email_val.lower().encode()).hexdigest()[:8]

    pptx_temp = os.path.join(output_dir, f"{safe_name}_{uid}_tmp.pptx")
    pdf_out = os.path.join(output_dir, f"{safe_name}_{uid}.pdf")

    try:
        _apply_replacements(template_path, pptx_temp, row)
        _convert_to_pdf(pptx_temp, pdf_out)
    finally:
        _delete_retry(pptx_temp)

    return pdf_out
